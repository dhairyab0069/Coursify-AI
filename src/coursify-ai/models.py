"""
Content Generation: PDFs, Slides, Quizzes
"""
import os
import random
import string
import tempfile
import time
from datetime import datetime
from io import BytesIO
from flask import jsonify, url_for
from flask_login import current_user
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation
from pptx.util import Pt
from docx import Document
from PyPDF2 import PdfReader
from bson.objectid import ObjectId

import extensions
from utils import llm, sanitize_filename, wrap_text, split_content_into_chunks, is_latex, render_latex, get_difficulty_text


# ============================================================================
# PDF GENERATION
# ============================================================================

def generate_pdf(prompt, length, difficulty):
    """Generate a PDF based on user input"""
    pdf_directory = os.path.join(os.getcwd(), 'pdfs')
    if not os.path.exists(pdf_directory):
        os.makedirs(pdf_directory)

    pdf_basename = sanitize_filename(prompt) or 'generated_file'
    timestamp = datetime.now().strftime('%Y-%m-%d')
    pdf_filename = f"{pdf_basename}_{timestamp}.pdf"
    diff = get_difficulty_text(difficulty)

    # Generate table of contents
    toc_prompt = f"Give Table of contents for topic: {prompt} (max 5 topics, each with 2 subtopics. Topics UPPERCASE, subtopics normal case. No 'Table of contents' heading, just contents. Difficulty: {diff})"
    toc = llm.generate(toc_prompt)

    if toc is None:
        return jsonify(success=False, error="Failed to generate content from LLM")

    # PDF setup
    font_name = "Helvetica"
    font_size = 12
    left_margin = 72
    top_margin = 720
    line_height = 14
    page_width, page_height = letter
    bottom_margin = 72
    content_width = page_width - 2 * left_margin

    pdf_path = os.path.join(pdf_directory, pdf_filename)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    c.setFont(font_name, font_size)

    # Draw Table of Contents
    y_position = top_margin
    c.drawString(left_margin, y_position, "Table of Contents")
    y_position -= line_height * 2

    for line in toc.split('\n'):
        if not line.strip():
            continue

        is_topic = line.isupper()

        if is_topic:
            c.setFont("Helvetica-Bold", 14)
            y_position -= 10
        else:
            c.setFont("Helvetica", 12)
            line = "   " + line

        c.drawString(left_margin, y_position, line)
        y_position -= line_height

        if y_position < bottom_margin:
            c.showPage()
            c.setFont("Helvetica", 12)
            y_position = top_margin

    # New page for content
    c.showPage()
    c.setFont("Helvetica", 12)
    y_position = top_margin

    lines = [l for l in toc.split('\n') if l.strip()]

    # Generate content for each line
    for line in lines:
        is_topic = line.isupper()

        if is_topic:
            c.setFont("Helvetica-Bold", 14)
            wrapped_text = wrap_text(line, content_width, font_name, 14)
        else:
            c.setFont("Helvetica", 12)

            # Generate content for subtopic
            content_prompt = f"Explain in 3-4 sentences: {line.strip()} (related to {prompt}, difficulty: {diff})"
            generated_content = llm.generate(content_prompt)

            if generated_content is None:
                return jsonify(success=False, error="Failed to generate content from LLM")

            # Write subtopic title
            c.drawString(left_margin, y_position, line)
            y_position -= line_height

            wrapped_text = wrap_text(generated_content, content_width, font_name, 12)

        # Write wrapped text
        for text_line in wrapped_text:
            if y_position < bottom_margin:
                c.showPage()
                c.setFont(font_name, 12)
                y_position = top_margin
            c.drawString(left_margin, y_position, text_line)
            y_position -= line_height

    c.save()
    print(f"PDF saved to {pdf_path}")

    # Save to GridFS
    if current_user.is_authenticated:
        with open(pdf_path, 'rb') as pdf_file:
            extensions.fs.put(pdf_file, filename=pdf_filename, user_id=current_user.user_id)

    pdf_url = url_for('get_pdf', filename=pdf_filename)
    return jsonify(success=True, pdf_url=pdf_url)


# ============================================================================
# SLIDES GENERATION
# ============================================================================

def generate_slides(prompt, length, difficulty):
    """Generate a presentation based on user input"""
    pptx_directory = os.path.join(os.getcwd(), 'presentations')
    if not os.path.exists(pptx_directory):
        os.makedirs(pptx_directory)

    pptx_basename = sanitize_filename(prompt) or 'generated_presentation'
    timestamp = datetime.now().strftime('%Y-%m-%d')
    pptx_filename = f"{pptx_basename}_{timestamp}.pptx"
    diff = get_difficulty_text(difficulty)

    # Generate table of contents
    toc_prompt = f"Give Table of contents for topic: {prompt} (Difficulty: {diff}, Length: {length} words total)"
    toc = llm.generate(toc_prompt)

    if toc is None:
        return jsonify(success=False, error="Failed to generate content from LLM")

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Table of Contents"

    # Add ToC to slide
    content = slide.placeholders[1]
    tf = content.text_frame
    for line in toc.split('\n'):
        if line.strip():
            p = tf.add_paragraph()
            p.text = line
            p.level = 0 if line.isupper() else 1

    # Generate content for each section
    toc_lines = [l.strip() for l in toc.split('\n') if l.strip()]
    words_per_section = length // max(len(toc_lines), 1)

    for line in toc_lines:
        if not line:
            continue

        section_prompt = f"Explain {line} in detail (Difficulty: {diff}, Length: {words_per_section} words)"
        section_content = llm.generate(section_prompt)

        if section_content:
            content_chunks = split_content_into_chunks(section_content)

            for chunk in content_chunks:
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = line

                content_box = slide.placeholders[1]
                tf = content_box.text_frame
                tf.text = chunk

                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)

    pptx_path = os.path.join(pptx_directory, pptx_filename)
    prs.save(pptx_path)

    # Save to GridFS
    if current_user.is_authenticated:
        with open(pptx_path, 'rb') as pptx_file:
            extensions.fs.put(pptx_file, filename=pptx_filename, user_id=current_user.user_id)

    pptx_url = url_for('get_presentation', filename=pptx_filename)
    return jsonify(success=True, pptx_url=pptx_url)


# ============================================================================
# QUIZ GENERATION
# ============================================================================

def generate_quiz_content(topic, subject, mcqs, true_false, short_questions, long_questions):
    """Generate quiz content"""
    mcqs = int(mcqs)
    true_false = int(true_false)
    short_questions = int(short_questions)
    long_questions = int(long_questions)

    quiz_content = f"Quiz on {topic} - {subject}\n\n"

    # Multiple Choice Questions
    if mcqs > 0:
        quiz_content += "Multiple Choice Questions:\n"
        for i in range(1, mcqs + 1):
            prompt = f"Generate a multiple choice question about {topic} and {subject} with 4 options (A, B, C, D) and mark the correct answer."
            mcq = llm.generate(prompt)
            if mcq is None:
                return None
            quiz_content += f"{i}. {mcq}\n\n"

    # True/False Questions
    if true_false > 0:
        quiz_content += "True/False Questions:\n"
        for i in range(1, true_false + 1):
            prompt = f"Generate a true or false question about {topic} and {subject}."
            tf_question = llm.generate(prompt)
            if tf_question is None:
                return None
            quiz_content += f"{i}. {tf_question}\n\n"

    # Short Answer Questions
    if short_questions > 0:
        quiz_content += "Short Answer Questions:\n"
        for i in range(1, short_questions + 1):
            prompt = f"Generate a short answer question about {topic} and {subject}."
            short_q = llm.generate(prompt)
            if short_q is None:
                return None
            quiz_content += f"{i}. {short_q}\n\n"

    # Long Answer Questions
    if long_questions > 0:
        quiz_content += "Long Answer Questions:\n"
        for i in range(1, long_questions + 1):
            prompt = f"Generate a detailed essay question about {topic} and {subject}."
            long_q = llm.generate(prompt)
            if long_q is None:
                return None
            quiz_content += f"{i}. {long_q}\n\n"

    return quiz_content


def generate_quiz_from_form(topic, subject, mcqs, true_false, short_questions, long_questions):
    """Generate quiz document from form data"""
    quiz_content = generate_quiz_content(topic, subject, mcqs, true_false, short_questions, long_questions)

    if quiz_content is None:
        return jsonify(success=False, error="Failed to generate quiz. Check LLM configuration.")

    doc = Document()
    doc.add_heading(f"Quiz on {topic} - {subject}", level=1)
    doc.add_paragraph(quiz_content)
    quiz_filename = f"{sanitize_filename(topic)}_{sanitize_filename(subject)}_quiz.docx"

    if not os.path.exists('quiz_files'):
        os.makedirs('quiz_files')
    quiz_path = os.path.join('quiz_files', quiz_filename)
    doc.save(quiz_path)

    with open(quiz_path, 'rb') as f:
        file_id = extensions.fs.put(f, filename=quiz_filename, user_id=current_user.get_id())

    quiz_url = url_for('get_doc', file_id=str(file_id), _external=True)
    return jsonify(success=True, quiz_url=quiz_url)


def generate_quiz_from_file(file_id):
    """Generate quiz from uploaded PDF/PPTX"""
    try:
        file_id = ObjectId(file_id)
    except:
        return jsonify({'error': 'Invalid file ID'}), 400

    try:
        file = extensions.fs.get(file_id)
    except:
        return jsonify({'error': 'File not found'}), 404

    file_name = file.filename
    file_stream = BytesIO(file.read())

    # Extract text based on file type
    file_extension = os.path.splitext(file.filename)[1]
    if file_extension == '.pdf':
        extracted_text = extract_text_from_pdf(file_stream)
    elif file_extension == '.pptx':
        extracted_text = extract_text_from_pptx(file_stream)
    else:
        return jsonify({'error': 'Unsupported file format'}), 400

    questions = generate_questions_from_text(extracted_text)

    doc = Document()
    doc.add_heading('Quiz', level=1)
    for i, question in enumerate(questions, start=1):
        doc.add_paragraph(f"Q{i}: {question}")

    temp_dir = tempfile.mkdtemp()
    timestamp = str(int(time.time()))
    random_string = ''.join(random.choices(string.ascii_lowercase + string.digits, k=8))
    doc_filename = file_name.split('.')[0] + '_quiz_' + timestamp + '_' + random_string + '.docx'

    temp_path = os.path.join(temp_dir, doc_filename)
    doc.save(temp_path)

    with open(temp_path, 'rb') as doc_file:
        if current_user.is_authenticated:
            file_id = extensions.fs.put(doc_file, filename=doc_filename, user_id=str(current_user.get_id()))

    os.remove(temp_path)
    os.rmdir(temp_dir)

    quiz_url = url_for('get_doc', file_id=file_id, _external=True)
    return jsonify(success=True, quiz_url=quiz_url)


def extract_text_from_pdf(file_stream):
    """Extract text from PDF file stream"""
    file_stream.seek(0)
    reader = PdfReader(file_stream)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text


def extract_text_from_pptx(file_stream):
    """Extract text from PPTX file stream"""
    file_stream.seek(0)
    presentation = Presentation(file_stream)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def generate_questions_from_text(text):
    """Generate quiz questions from text"""
    mcq_prompt = f"Create 10 multiple-choice questions based on: \n{text[:2000]}\nEach with 4 options (A, B, C, D) and correct answer."
    short_ans_prompt = f"Create 3 short answer questions based on: \n{text[:2000]}"

    mcq_text = llm.generate(mcq_prompt)
    short_ans_text = llm.generate(short_ans_prompt)

    questions = []
    if mcq_text:
        questions.extend(mcq_text.strip().split('\n\n'))
    if short_ans_text:
        questions.extend(short_ans_text.strip().split('\n\n'))

    return questions if questions else ["Failed to generate questions."]