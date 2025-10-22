import base64
import glob
import io
import logging
from mailbox import Message
from pydoc_data import topics
import subprocess
import time
from urllib.parse import unquote_to_bytes
from flask import after_this_request, flash, session
from datetime import datetime
import os
import re
import random
import string
from click import wrap_text
from flask import Flask, jsonify, render_template, request, send_from_directory, url_for, Response, send_file, make_response, redirect
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase.pdfmetrics import stringWidth
from flask_cors import CORS
import openai 
import matplotlib.pyplot as plt
from io import BytesIO
from pymongo import MongoClient
from gridfs import GridFS, NoFile
from PyPDF2 import PdfReader
from werkzeug.utils import secure_filename
from bson.objectid import ObjectId

from flask_bcrypt import Bcrypt
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user
from pymongo import MongoClient
from bson.objectid import ObjectId
from itsdangerous import URLSafeTimedSerializer, SignatureExpired, BadSignature
from sendgrid import SendGridAPIClient
from sendgrid.helpers.mail import Mail
from flask_mail import Mail, Message
from pptx import Presentation
from flask import send_file, abort
from flask_wtf import FlaskForm
from wtforms import StringField, PasswordField, SubmitField, SelectField, RadioField, TextAreaField
from wtforms.validators import DataRequired, Email, EqualTo, Length, ValidationError
from datetime import datetime
from docx import Document
import tempfile
import subprocess
import os
from PIL import Image
from urllib.parse import quote
import shutil
from pathlib import Path
from pptx.util import Pt

import subprocess
from collections import defaultdict
from collections import Counter



class RegistrationForm(FlaskForm):
    '''Registration form class. Inherits from FlaskForm. '''
    first_name = StringField('First Name', validators=[DataRequired()])
    last_name = StringField('Last Name', validators=[DataRequired()])
    email = StringField('Email', validators=[DataRequired(), Email()])
    password = PasswordField('Password', validators=[
        DataRequired(),
        Length(min=8, message='Password must be at least 8 characters long.'),
        EqualTo('confirm_password', message='Passwords must match.')
    ])
    confirm_password = PasswordField('Confirm Password')
    birth_date = SelectField('Birth Date', choices=[(str(i), str(i)) for i in range(1, 32)])
    birth_month = SelectField('Birth Month', choices=[(str(i), str(i)) for i in range(1, 13)])
    birth_year = SelectField('Birth Year', choices=[(str(i), str(i)) for i in range(1900, 2022)])
    gender = RadioField('Gender', choices=[('Male', 'Male'), ('Female', 'Female'), ('Custom', 'Custom')])
    custom_gender = TextAreaField('Custom Gender')
    submit = SubmitField('Register')



app = Flask(__name__, template_folder='my_templates')
CORS(app)
bcrypt = Bcrypt(app)
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'
openai.api_key = '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'


app.secret_key = 'coursifyai1234'

serializer = URLSafeTimedSerializer(app.secret_key)




# Setup MongoDB connection
client = MongoClient('mongodb+srv://Remy:1234@cluster0.vgzdbrr.mongodb.net/')
db = client['new_pdfs']
db2=client['Login_details']
db3=client['reviews']
fs = GridFS(db)
users_collection=db2.users
reviews_collection=db3.reviews

app.config['MAIL_SERVER'] = 'smtp-mail.outlook.com'
app.config['MAIL_PORT'] = 587
app.config['MAIL_USE_TLS'] = True
app.config['MAIL_USERNAME'] = 'coursify@outlook.com'  
app.config['MAIL_PASSWORD'] = 'Gunners4Eva$'

mail = Mail(app)



# User model
class User(UserMixin):
    '''User class for Flask-Login.'''
    def __init__(self, user_id, email):
        self.user_id = str(user_id)
        self.email = email

    def get_id(self):
        return self.user_id

# User Loader
@login_manager.user_loader
def load_user(user_id):
    '''Function to load a user from the database.'''
    user = users_collection.find_one({"_id": ObjectId(user_id)})
    if not user:
        return None
    return User(user_id=user["_id"], email=user["email"])

@app.route('/')
def home():
    '''Function to render the homepage.'''
    # If the user is authenticated, redirect to the dashboard
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    # Otherwise, show the homepage with login and register options
    return render_template('homepage.html')

SENDGRID_API_KEY = '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'  

def validate_password(password):
    '''Validation of password'''
    if len(password) < 8:
        return "Password must be at least 8 characters long."
    if not re.search("[a-z]", password):
        return "Password must contain at least one lowercase letter."
    if not re.search("[A-Z]", password):
        return "Password must contain at least one uppercase letter."
    if not re.search("[0-9]", password):
        return "Password must contain at least one number."
    if not re.search("[!@#$%^&*]", password):
        return "Password must contain at least one special character (!@#$%^&*)."
    return None
   
# Registration Route
@app.route('/register', methods=['GET', 'POST'])
def register():
    '''Function to register a new user.'''
    if request.method == 'POST':
        first_name = request.form.get('first_name')
        last_name = request.form.get('last_name')
        email = request.form.get('email')
        password = request.form.get('password')
        birth_day = int(request.form.get('birth_day'))
        birth_month = int(request.form.get('birth_month'))
        birth_year = int(request.form.get('birth_year'))
        gender = request.form.get('gender')
        custom_gender = request.form.get('custom_gender') if 'custom_gender' in request.form else None
        
        if gender == 'custom' and custom_gender:
            gender_value = custom_gender
        else:
            gender_value = gender
        
        # Validate if user exists
        user = users_collection.find_one({"email": email})
        if user:
            flash('Email already exists.')
            return redirect(url_for('register'))
        
        # Validate the password
        password_error = validate_password(password)
        if password_error:
            flash(password_error)
            return redirect(url_for('register'))
        
        hashed_password = bcrypt.generate_password_hash(password).decode('utf-8')
        
        # Construct birthdate and gender
        birth_date = datetime(birth_year, birth_month, birth_day)
        
        # Insert new user into MongoDB with 'verified' set to False initially
        user_id = users_collection.insert_one({
            "first_name": first_name,
            "last_name": last_name,
            "email": email,
            "password": hashed_password,
            "birth_date": birth_date,
            "gender": gender_value,
            "verified": False
        }).inserted_id
        
        # Generate a token for email verification
        token = serializer.dumps(email, salt='email-confirm')
        
        # Construct the URL for email verification
        confirm_url = url_for('confirm_email', token=token, _external=True)
        
        # Prepare the verification email content
        html_content = render_template('email_verification.html', confirm_url=confirm_url)
        subject = "Please confirm your email"
        
        # Use the existing function to send the email
        send_email(email, subject, html_content)

        flash('A confirmation email has been sent. Please check your inbox.', 'info')
        return redirect(url_for('login'))
    
    return render_template('register.html')

# Send Email Function
def send_email(to_email, subject, html_content):
    '''Function to send an email to the user for verification.'''
    msg = Message(subject,
                  sender=app.config['MAIL_USERNAME'],
                  recipients=[to_email],
                  html=html_content)
    print(html_content)
    mail.send(msg)


# Email Confirmation Route
@app.route('/confirm/<token>')
def confirm_email(token):
    try:
        email = serializer.loads(token, salt='email-confirm', max_age=3600)  # 1 hour to verify
    except (SignatureExpired, BadSignature):
        return 'The confirmation link is invalid or has expired.', 400

    user = users_collection.find_one({"email": email})
    if user and not user.get('verified'):
        users_collection.update_one({'_id': user['_id']}, {'$set': {'verified': True}})
        flash('Your email has been verified!', 'success')
    else:
        flash('Your email is already verified or the user does not exist.', 'warning')
    return redirect(url_for('login'))

# Login Route
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        email = request.form['email']
        password = request.form['password']

        user = users_collection.find_one({"email": email})
        if user and bcrypt.check_password_hash(user['password'], password):
            if not user.get('verified', False):
                # User found but not verified
                flash('Please verify your email address before logging in.', 'warning')
                return redirect(url_for('login'))

            # User is verified, proceed to login
            user_obj = User(user_id=str(user["_id"]), email=email)
            login_user(user_obj, remember='remember' in request.form)
            return redirect(url_for('index'))
        else:
            flash('Invalid credentials. Please try again.', 'danger')

    return render_template('login.html')


@app.route('/index')
@login_required
def index():
    '''Function to render the dashboard page after successful login.'''
    return render_template('index.html')

@app.route('/logout')
def logout():
    '''Function to log out the user.'''
    logout_user()
    return redirect(url_for('home'))

@app.route('/quiz_generate')
def quiz_generate():
    '''Function to render the quiz generation page.'''
    return render_template('quiz.html')


# SETINGS PAGE
@app.route('/settings.html')
@login_required
def settings_html():
    '''Function to render the settings page.'''
    user_id = current_user.get_id()
    user = users_collection.find_one({"_id": ObjectId(user_id)})
    if user:
        return render_template('settings.html', user=user)
    else:
        flash("User not found.")
        return redirect(url_for('index'))
 
#ACCOUNT SETTINGS - UPDATE (FIRST / LAST NAME, EMAIL)
@app.route('/update_account_settings', methods=['POST'])
@login_required
def update_account_settings():
    '''Function to update the user's account settings.'''
    first_name = request.form.get('firstname')
    last_name = request.form.get('lastname')
    email = request.form.get('email')
    
    user_id = current_user.get_id()
    user = users_collection.find_one({"_id": ObjectId(user_id)})

    if user:
        updates = {}
        if user.get('first_name') != first_name:
            updates['first_name'] = first_name

        if user.get('last_name') != last_name:
            updates['last_name'] = last_name

        if user.get('email') != email:
            updates['email'] = email

        if updates:
            users_collection.update_one({"_id":ObjectId(user_id)}, {"$set":updates})
            flash('Your account has been updated successfully')
        else:
            flash('No changes were made to your account.')
    else:
        flash('User not found.')
        return redirect(url_for('settings_html'))

    return redirect(url_for('settings_html'))

@app.route('/change_password', methods=['POST'])
@login_required
def change_password():
    '''Function to change the user's password.'''
    current_password = request.form.get('current-password')
    new_password = request.form.get('new-password')
    confirm_new_password = request.form.get('confirm-new-password')

    user_id = current_user.get_id()
    user = users_collection.find_one({"_id": ObjectId(user_id)})

    if not user:
        return redirect(url_for('settings_html'))
    
    if not bcrypt.check_password_hash(user['password'], current_password): 
        flash('Current password is incorrect.','settings_page')
        return redirect(url_for('settings_html'))

    if new_password != confirm_new_password:
        flash('New passwords do not match.','settings_page')
        return redirect(url_for('settings_html'))

    hashed_new_password = bcrypt.generate_password_hash(new_password).decode('utf-8')
    users_collection.update_one({"_id": ObjectId(user_id)}, {"$set": {"password": hashed_new_password}})

    flash('Your password has been updated successfully.','settings_page')
    return redirect(url_for('settings_html'))

def send_pw_reset_email(recipient_email, reset_url):
    subject = "Password Reset Request"
    sender_email = 'coursify@outlook.com'
    recipients = [recipient_email]
    
    msg = Message(subject,
                  sender=sender_email,
                  recipients=recipients)
   
    msg.html = f"<p>Please click on the following link to reset your password:</p><a href='{reset_url}'>{reset_url}</a>"
    mail.send(msg)

@app.route('/quiz-generator-page', methods=['POST'])
def quiz_generator_page():
    '''Function to generate a quiz based on the user's input.'''
    topic = request.form.get('topic')
    subject = request.form.get('subject')
    mcqs = request.form.get('mcqs')
    true_false = request.form.get('true-false')
    short_questions = request.form.get('short-questions')
    long_questions = request.form.get('long-questions')

    quiz_content = quiz_page_generate(topic, subject, mcqs, true_false, short_questions, long_questions)

    doc = Document()
    doc.add_heading(f"Quiz on {topic} - {subject}", level=1)
    doc.add_paragraph(quiz_content)
    quiz_filename = f"{topic}_{subject}_quiz.docx"
    
    if not os.path.exists('quiz_files'):
        print("Creating directory at:", os.path.abspath('quiz_files'))
        os.makedirs('quiz_files')
    quiz_path = os.path.join('quiz_files', quiz_filename)
    doc.save(quiz_path)

    with open(quiz_path, 'rb') as f:
        file_id = fs.put(f, filename=quiz_filename, user_id=current_user.get_id())

    quiz_url = url_for('get_doc', file_id=str(file_id), _external=True)
    return jsonify(success=True, quiz_url=quiz_url)


def quiz_page_generate(topic, subject, mcqs, true_false, short_questions, long_questions):
    '''Function to generate a quiz based on the user's input.'''
    mcqs = int(mcqs)
    true_false = int(true_false)
    short_questions = int(short_questions)
    long_questions = int(long_questions)

    quiz_content = f"Quiz on {topic} - {subject}\n\n"

    quiz_content += "Multiple Choice Questions:\n"
    for i in range(1, mcqs + 1):
        prompt = f"Generate a multiple choice question related to {topic} and {subject}."
        mcq = call_openai_api(prompt)
        if mcq is None:
            return jsonify(success=False, error="Failed to generate text from OpenAI API")
        quiz_content += f"{i}. {mcq}\n\n"

    quiz_content += "True/False Questions:\n"
    for i in range(1, true_false + 1):
        prompt = f"Generate a true or false question related to {topic} and {subject}."
        true_false_question = call_openai_api(prompt)
        if true_false_question is None:
            return jsonify(success=False, error="Failed to generate text from OpenAI API")
        quiz_content += f"{i}. {true_false_question}\n\n"

    quiz_content += "Short Answer Questions:\n"
    for i in range(1, short_questions + 1):
        prompt = f"Generate a short answer question related to {topic} and {subject}."
        short_question = call_openai_api(prompt)
        if short_question is None:
            return jsonify(success=False, error="Failed to generate text from OpenAI API")
        quiz_content += f"{i}. {short_question}\n\n"

    quiz_content += "Long Answer Questions:\n"
    for i in range(1, long_questions + 1):
        prompt = f"Generate a long answer question related to {topic} and {subject}."
        long_question = call_openai_api(prompt)
        if long_question is None:
            return jsonify(success=False, error="Failed to generate text from OpenAI API")
        quiz_content += f"{i}. {long_question}\n\n"

    return quiz_content


@app.route('/forgot_password', methods=['GET','POST'])
def forgot_password():
    if request.method == 'POST':
        email = request.form.get('email')
        user = users_collection.find_one({"email": email}) 
        if user:
            token = serializer.dumps(email, salt='password-reset-salt')
            reset_url = url_for('reset_forgot_password', token=token, _external=True)
            send_pw_reset_email(email, reset_url)
            flash('An email has been sent with instructions to reset your password.', 'reset_password_sent')
        else:
            flash('User does not exist. Please check the email address.', 'reset_password_sent')

    return render_template('forgot_password.html')

@app.route('/reset_forgot_password/<token>', methods=['GET', 'POST'])
def reset_forgot_password(token):
    try:
        email = serializer.loads(token, salt='password-reset-salt', max_age=3600)
    except SignatureExpired:
        flash('The password reset link has expired.', 'reset_password_page')
        return redirect(url_for('forgot_password'))
    except BadSignature:
        flash('The password reset link is invalid.', 'reset_password_page')
        return redirect(url_for('forgot_password'))

    if request.method == 'POST':
        new_password = request.form.get('new_password')
        confirm_new_password = request.form.get('confirm_new_password')

        if new_password != confirm_new_password:
            flash('New passwords do not match.', 'reset_password_page')
            return redirect(url_for('reset_forgot_password', token=token))

        password_error = validate_password(new_password)
        if password_error:
            flash(password_error)
            return redirect(url_for('reset_forgot_password', token=token))

        hashed_new_password = bcrypt.generate_password_hash(new_password).decode('utf-8')
        users_collection.update_one({"email": email}, {"$set": {"password": hashed_new_password}})

        flash('Your password has been reset successfully.', 'login_page')
        return redirect(url_for('login'))

    return render_template('reset_password.html', token=token)

@app.route('/share_via_email', methods=['POST'])
def share_via_email():
    '''Function to share a file via email as an attachment.'''
    recipient = request.form.get('email')
    file_id = request.form.get('file_id')

    file_obj = fs.get(ObjectId(file_id))
    file_data = file_obj.read()
    file_name = file_obj.filename

    msg = Message('Your Shared File',
                  sender='coursify@outlook.com',
                  recipients=[recipient])
    msg.body = 'Please find the attached file.'
    msg.attach(file_name, "application/octet-stream", file_data)
    mail.send(msg)

    flash('Email sent with the file attached!')
    return redirect(url_for('my_content'))

        
@app.route('/share/<file_id>')
def share_file(file_id):
    '''Function to share a file with a unique URL.'''
    file_id = ObjectId(file_id)
    file = fs.get(file_id)
    response = make_response(file.read())
    response.mimetype = 'application/pdf'
    response.headers.set('Content-Disposition', 'attachment', filename=file.filename)
    return response


def extract_text_from_pdf(pdf_path):
    '''Extract text from a PDF using PdfReader.'''
    with open(pdf_path, 'rb') as f:
        reader = PdfReader(f)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
    return text

@app.route('/content')
def my_content():
    '''Function to display the content page.'''
    if current_user.is_authenticated:
        fs=GridFS(db)
        files = fs.find({"user_id": current_user.get_id()}).sort("_id",-1)
        file_data = []
        for file in files:
            file_data.append({"filename": file.filename, "_id": str(file._id)})
        return render_template('content.html', file_data=file_data)

@app.route('/ai.html')
def ai_html():
    return render_template('ai.html')

@app.route('/preview.html')
def preview():
    return render_template('preview.html')

@app.route('/faq.html')
def faq():
    return render_template('faq.html')

@app.route('/api/chat', methods=['POST'])
def chat():
    '''Function for the chatbot API endpoint.'''
    try:
        user_input = request.json['message']
        print("User input received:", user_input)
        
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "AI, your role is to assist users in navigating and utilizing the features of Coursify.ai effectively. You will only answer questions related to this website. When a user inquires about generating content, guide them through entering the desired content length, difficulty, and topics. If they need to manage their notes, direct them to the 'My Content' section. For account-related queries, help users find where to update their details or change their password on the 'Settings' page. Remember to always respect user privacy and handle their data with care. Should users need assistance with features or encounter issues, use the AI assistant chat to offer support. Provide clear, friendly guidance, ensuring users feel supported at every step. For new users, explain the registration process, and for returning users, assist them with the login procedure. Always communicate in a warm, professional manner, creating a positive user experience."},
                {"role": "user", "content": user_input}
            ]
        )
        print("Response from OpenAI:", response)

        ai_reply = response['choices'][0]['message']['content']
        print("AI Reply:", ai_reply)

        return jsonify({'reply': ai_reply})
    except Exception as e:
        print("An error occurred:", e)
        return jsonify({'error': str(e)}), 500

def is_latex(text):
    '''Check if a string contains LaTeX content.'''
    return bool(re.search(r"\$.*\$", text))

def render_latex(formula, fontsize=12, dpi=150):
    '''Render a LaTeX formula to an image and return the image as a BytesIO buffer.'''
    plt.rcParams['text.usetex'] = True
    fig = plt.figure()
    fig.patch.set_alpha(0)
    ax = fig.add_subplot(111)
    ax.axis('off')
    ax.patch.set_alpha(0)
    ax.text(0, 0, f"\\begin{{equation}}{formula}\\end{{equation}}", fontsize=fontsize)
    buffer = BytesIO()
    fig.savefig(buffer, dpi=dpi, transparent=True, format='png')
    plt.close(fig)
    buffer.seek(0)
    return buffer


def sanitize_filename(text):
    '''Sanitize a string to be used as a filename.'''
    valid_chars = "-_.() %s%s" % (string.ascii_letters, string.digits)
    sanitized = ''.join(c for c in text if c in valid_chars)
    return sanitized[:255]


def wrap_text(text, max_width, font_name, font_size):
    '''Wrap text to fit within a given width.'''
    wrapped_text = []
    words = text.split()

    while words:
        line = ''
        while words and stringWidth(line + words[0], font_name, font_size) <= max_width:
            line += (words.pop(0) + ' ')
        wrapped_text.append(line)

    return wrapped_text

def content(prompt, length):
    '''Calls the OpenAI API to generate content based on the user's input.'''
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a content generation tool. Give as much accurate information as you are supposed to give and follow the instruction in the prompt."},
                {"role": "user", "content": prompt}
            ]
        )
        return response['choices'][0]['message']['content']
    except Exception as e:
        print(f"Error calling OpenAI API: {e}")
        return None

def call_openai_api(prompt):
    '''Calls the OpenAI API to generate content based on the user's input.'''
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a content generation tool. Give as much accurate information as you are supposed to give and follow the instruction in the prompt."},
                {"role": "user", "content": prompt}
            ]
        )
        return response['choices'][0]['message']['content']
    except Exception as e:
        print(f"Error calling OpenAI API: {e}")
        return None
    

@app.route('/generate_content', methods=['POST'])
def generate_content():
    '''Function to generate content based on the user's input.'''
    length = request.form.get('length', type=int, default=100)
    prompt = request.form.get('topics', default='')
    difficulty = request.form.get('difficulty', type=int, default=1)
    content_format = request.form.get('format', default='pdf')

    if content_format == 'pdf':
        return generate_pdf(prompt, length, difficulty)
    elif content_format == 'slides':
        return generate_slides(prompt, length, difficulty)
    else:
        return jsonify(success=False, error="Invalid format selected")

def generate_pdf(prompt, length, difficulty):
    '''Generate a PDF based on the user's input.'''
    pdf_directory = os.path.join(os.getcwd(), 'pdfs')
    if not os.path.exists(pdf_directory):
        os.makedirs(pdf_directory)

    pdf_basename = sanitize_filename(prompt)
    pdf_basename = pdf_basename if pdf_basename else 'generated_file'

    timestamp = datetime.now().strftime('%Y-%m-%d')
    unique_suffix = ''.join(random.choices(string.ascii_letters + string.digits, k=6))
    pdf_filename = f"{pdf_basename}_{timestamp}.pdf"

    if difficulty == 1:
        diff = "Grade 1 - Introduction to fundamental concepts"
    elif difficulty == 2:
        diff = "Grade 2 - Building on basic skills and knowledge"
    elif difficulty == 3:
        diff = "Grade 3 - Expanding core understanding"
    elif difficulty == 4:
        diff = "Grade 4 - Deepening comprehension of key topics"
    elif difficulty == 5:
        diff = "Grade 5 - Preparing for intermediate challenges"
    elif difficulty == 6:
        diff = "Grade 6 - Transitioning to more complex subjects"
    elif difficulty == 7:
        diff = "Grade 7 - Enhancing critical thinking and application"
    elif difficulty == 8:
        diff = "Grade 8 - Solidifying foundational knowledge"
    elif difficulty == 9:
        diff = "Grade 9 - High school introductory concepts"
    elif difficulty == 10:
        diff = "Grade 10 - Sophomore explorations and depth"
    elif difficulty == 11:
        diff = "Grade 11 - Junior year, college prep, and advanced topics"
    elif difficulty == 12:
        diff = "Grade 12 - Senior year, culmination, and readiness for next steps"

    toc = call_openai_api("Give Table of contents for topic: " +prompt+"(such that there are max 5 topics and each topic has two sub topics. Topics should be upper Case and subtopics otherwise. Dont Start with heading : Table of contents, just show contents with difficulty " + diff)
    if toc is None:
        return jsonify(success=False, error="Failed to generate text from OpenAI API")

    font_name = "Helvetica"
    font_size = 12
    left_margin = 72
    top_margin = 720
    line_height = 14
    page_width, page_height = letter
    bottom_margin = 72
    content_width = page_width - 2 * left_margin

    pdf_path = os.path.join(pdf_directory, pdf_filename)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    c.setFont(font_name, font_size)
    toc_dict = {}
    current_topic = ""

    y_position = top_margin
    c.drawString(left_margin, y_position, "Table of Contents")
    y_position -= line_height
    for line in toc.split('\n'):
        is_topic = line.isupper()
        contains_latex = is_latex(line)

        if contains_latex:
            latex_image_io = render_latex(line[1:-1])
            c.drawImage(latex_image_io, left_margin, y_position, width=200, height=100, preserveAspectRatio=True, anchor='n')
            y_position -= 100
        else:
            if is_topic:
                c.setFont("Helvetica-Bold", 14)
                y_position -= 20
            else:
                if current_topic:
                    toc_dict[current_topic].append(line.strip())
                c.setFont("Helvetica", 12)
                line = "   " + line

            c.drawString(left_margin, y_position, line)
            y_position -= line_height

        if y_position < 100:
            c.showPage()
            c.setFont("Helvetica", 12)
            y_position = top_margin
    
    c.showPage()
    c.setFont("Helvetica", 12)
    y_position = top_margin

    lines = toc.split('\n')
    
    print(lines)
    for line in lines:
        is_topic = line.isupper()
        contains_latex = is_latex(line)

        if contains_latex:
            latex_image_io = render_latex(line[1:-1])
            c.drawImage(latex_image_io, left_margin, y_position, width=200, height=100, preserveAspectRatio=True, anchor='n')
            y_position -= 100
        else:
            if is_topic:
                c.setFont("Helvetica-Bold", 14)
                wrapped_text = wrap_text(line, content_width, font_name, 14)
                print(wrapped_text)
            else:
                c.setFont("Helvetica", 12)
                c.drawString(left_margin, y_position, line)
                y_position -= line_height 
                prompt2 = content("explain in 3 lines about the following topic " + line + " related to " + prompt, "length : " + str(length/(len(lines))) + "words and difficulty level is " + diff)

                if prompt2 is None:
                    return jsonify(success=False, error="Failed to generate text from OpenAI API")
                wrapped_text = wrap_text(prompt2, content_width, font_name, 12)
                print ("sub topic done")
                
            for text_line in wrapped_text:
                if y_position < bottom_margin:
                    c.showPage()
                    c.setFont(font_name, 12)
                    y_position = top_margin
                c.drawString(left_margin, y_position, text_line)
                y_position -= line_height 

    c.save()
    print("PDF saved to", pdf_path)

    if current_user.is_authenticated:
        fs = GridFS(db)
        with open(pdf_path, 'rb') as pdf_file:
            fs.put(pdf_file, filename=pdf_filename, user_id=current_user.user_id)

    pdf_url = url_for('get_pdf', filename=pdf_filename)
    return jsonify(success=True, pdf_url=pdf_url)

def generate_slides(prompt, length, difficulty):
    '''Generate a presentation based on the user's input.'''
    pptx_directory = os.path.join(os.getcwd(), 'presentations')
    if not os.path.exists(pptx_directory):
        os.makedirs(pptx_directory)

    if difficulty == 1:
        diff = "Grade 1 - Introduction to fundamental concepts"
    elif difficulty == 2:
        diff = "Grade 2 - Building on basic skills and knowledge"
    elif difficulty == 3:
        diff = "Grade 3 - Expanding core understanding"
    elif difficulty == 4:
        diff = "Grade 4 - Deepening comprehension of key topics"
    elif difficulty == 5:
        diff = "Grade 5 - Preparing for intermediate challenges"
    elif difficulty == 6:
        diff = "Grade 6 - Transitioning to more complex subjects"
    elif difficulty == 7:
        diff = "Grade 7 - Enhancing critical thinking and application"
    elif difficulty == 8:
        diff = "Grade 8 - Solidifying foundational knowledge"
    elif difficulty == 9:
        diff = "Grade 9 - High school introductory concepts"
    elif difficulty == 10:
        diff = "Grade 10 - Sophomore explorations and depth"
    elif difficulty == 11:
        diff = "Grade 11 - Junior year, college prep, and advanced topics"
    elif difficulty == 12:
        diff = "Grade 12 - Senior year, culmination, and readiness for next steps"

    pptx_basename = sanitize_filename(prompt)
    pptx_basename = pptx_basename if pptx_basename else 'generated_presentation'
    timestamp = datetime.now().strftime('%Y-%m-%d')
    unique_suffix = ''.join(random.choices(string.ascii_letters + string.digits, k=6))
    pptx_filename = f"{pptx_basename}_{timestamp}.pptx"

    toc = call_openai_api("Give Table of contents for topic: " + prompt + "..." + "Difficulty level is " + diff + "..." + "Length of content is " + str(length) + " words.")
    if toc is None:
        return jsonify(success=False, error="Failed to generate text from OpenAI API")

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Table of Contents"

    toc_chunks = split_content_into_chunks(toc)

    for chunk in toc_chunks:
        slide = prs.slides.add_slide(slide_layout)
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = chunk
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(15)

    content = slide.placeholders[1]
    tf = content.text_frame

    for line in toc.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.level = 0 if line.isupper() else 1

    length = length // len(toc.split('\n'))
    for line in toc.split('\n'):
        if line.strip():
            section_content = call_openai_api(f"Explain {line.strip()} in detail. Difficulty: " + diff + "..." + "Length of content is " + str(length) + " words.")
            content_chunks = split_content_into_chunks(section_content)

            for chunk in content_chunks:
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = line.strip()

                content_box = slide.placeholders[1]
                tf = content_box.text_frame
                tf.text = chunk
                print("Topic done")

    pptx_path = os.path.join(pptx_directory, pptx_filename)
    prs.save(pptx_path)

    if current_user.is_authenticated:
        fs = GridFS(db)
        with open(pptx_path, 'rb') as pptx_file:
            fs.put(pptx_file, filename=pptx_filename, user_id=current_user.user_id)

    pptx_url = url_for('get_presentation', filename=pptx_filename)
    return jsonify(success=True, pptx_url=pptx_url)

def split_content_into_chunks(content):
    '''Split content into chunks that fit on a slide.'''
    max_chars_per_chunk = 400
    content_chunks = []
    current_chunk = ''
    for line in content.split('\n'):
        if len(current_chunk) + len(line) <= max_chars_per_chunk:
            current_chunk += line + '\n'
        else:
            content_chunks.append(current_chunk)
            current_chunk = line + '\n'
    content_chunks.append(current_chunk)
    return content_chunks

@app.route('/get_user_pdfs', methods=['GET'])
def get_user_pdfs():
    '''Get all PDFs uploaded by the currently logged in user.'''
    if current_user.is_authenticated:
        fs = GridFS(db)
        user_pdfs = fs.find({'user_id': current_user.user_id})
    return user_pdfs

@app.route('/pdf/<filename>')
def get_pdf(filename):
    '''Download a PDF from GridFS based on its filename.'''
    if current_user.is_authenticated:
        fs = GridFS(db)
        grid_out = fs.find_one({'filename': filename, 'user_id': current_user.user_id})

        if grid_out:
            response = make_response(grid_out.read())
            response.mimetype = 'application/pdf'
            return response
        else:
            return "File not found", 404
    else:
        return "Unauthorized", 401

@app.route('/presentation/<filename>')
def get_presentation(filename):
    '''Download a presentation from GridFS based on its filename.'''
    if current_user.is_authenticated:
        fs = GridFS(db)
        grid_out = fs.find_one({'filename': filename, 'user_id': current_user.user_id})

        if grid_out:
            response = make_response(grid_out.read())
            response.mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            response.headers["Content-Disposition"] = f"attachment; filename={filename}"
            return response
        else:
            abort(404, description="Presentation file not found")
    else:
        abort(401, description="Unauthorized to access this presentation")

@app.route('/get_doc/<file_id>')
def get_doc(file_id):
    '''Download a file from GridFS based on its file_id.'''
    try:
        file_id = ObjectId(file_id)
        grid_out = fs.get(file_id)
        return send_file(grid_out, download_name=grid_out.filename, as_attachment=True)
    except NoFile:
        return jsonify({'error': 'File not found'}), 404

@app.route('/list_pdfs', methods=['GET'])
def list_pdfs():
    '''List all PDFs in the database.'''
    if current_user.is_authenticated:
        fs = GridFS(db)
        user_pdfs = fs.find({'user_id': current_user.user_id})
    return render_template('content.html', pdfs=user_pdfs)    
  
@app.route('/check_file/<filename>', methods=['GET'])
def check_file(filename):
    '''Check if a file exists in the database.'''
    try:
        file = fs.get_last_version(filename=filename)
        if file:
            return jsonify(success=True, message="File exists in the database.")
    except:
        return jsonify(success=False, message="File does not exist in the database.")
    
@app.route('/submit_review', methods=['POST'])
@login_required
def submit_review():
    '''Submit a review for a file.'''
    form_data = request.form
    star_rating = form_data.get('star_rating')
    review_text = form_data.get('review_text')
    subject = form_data.get('subject')
    file_id = form_data.get('file_id')

    user_id = ObjectId(current_user.get_id())
    user_details = users_collection.find_one({"_id": user_id})
    
    if not user_details:
        return redirect(url_for('my_content'))

    first_name = user_details.get('first_name', '')
    last_name = user_details.get('last_name', '')

    if reviews_collection.find_one({"user_id": user_id, "file_id": file_id}):
        return redirect(url_for('my_content'))

    review = {
        "user_id": user_id,
        "file_id": file_id,
        "first_name": first_name,
        "last_name": last_name,
        "star_rating": star_rating,
        "review_text": review_text,
        "subject": subject,
        "timestamp": datetime.utcnow()
    }

    reviews_collection.insert_one(review)
    return redirect(url_for('my_content'))

@app.route('/reviews')
@login_required
def reviews():
    '''Display all reviews.'''
    user_id = current_user.get_id()  
    all_reviews = reviews_collection.find().sort("timestamp", -1)
    return render_template('reviews.html', reviews=all_reviews)

@app.route('/check_review_existence', methods=['GET'])
@login_required
def check_review_existence():
    '''Check if a review exists for a file.'''
    file_id = request.args.get('file_id')
    user_id = ObjectId(current_user.get_id())

    existing_review = reviews_collection.find_one({"user_id": str(user_id), "file_id": file_id})
    if existing_review:
        return jsonify({"reviewExists": True})
    else:
        return jsonify({"reviewExists": False})

@app.route('/delete_review', methods=['POST'])
@login_required
def delete_review():
    '''Delete a review.'''
    file_id = request.form.get('file_id')
    user_id = current_user.get_id()
    result = reviews_collection.delete_one({"user_id": user_id, "file_id": file_id})
    return redirect(url_for('reviews'))

@app.route('/delete/<file_id>', methods=['POST'])
def delete_file(file_id):
    '''Delete a file from GridFS based on its file_id.'''
    file_id = ObjectId(file_id)
    fs.delete(file_id)
    return redirect(url_for('my_content'))

@app.route("/chatbot", methods=["POST"])
def chatbot():
    '''Chatbot route to handle user messages and generate responses.'''
    message = request.json.get("message")
    openai.api_key = 'your-api-key'

    completion = openai.Completion.create(
        engine="text-davinci-003",
        prompt=message,
        max_tokens=150
    )

    if completion.choices and completion.choices[0].text.strip() != "":
        return completion.choices[0].text.strip()
    else:
        return 'Failed to Generate response!'

@app.route('/generate_quiz/<file_id>')
def generate_quiz(file_id):
    '''Generate a quiz based on the text extracted from a PDF or PPTX file.'''
    try:
        file_id = ObjectId(file_id)
    except:
        return jsonify({'error': 'Invalid file ID'}), 400

    try:
        file = fs.get(file_id)
    except:
        return jsonify({'error': 'File not found'}), 404

    file_name = file.filename
    file_stream = io.BytesIO(file.read())

    file_extension = os.path.splitext(file.filename)[1]
    if file_extension == '.pdf':
        extracted_text = extract_text_from_pdf(file_stream)
    elif file_extension == '.pptx':
        extracted_text = extract_text_from_pptx(file_stream)
    else:
        return jsonify({'error': 'Unsupported file format'}), 400

    questions = generate_questions(extracted_text)

    doc = Document()
    doc.add_heading('Quiz', level=1)
    for i, question in enumerate(questions, start=1):
        doc.add_paragraph(f"Q{i}: {question}")

    temp_dir = tempfile.mkdtemp()
    timestamp = str(int(time.time()))
    random_string = ''.join(random.choices(string.ascii_lowercase + string.digits, k=8))
    doc_filename = file_name.split('.')[0] + '_quiz_' + timestamp + '_' + random_string + '.docx'
    
    temp_path = os.path.join(temp_dir, doc_filename)
    doc.save(temp_path)

    with open(temp_path, 'rb') as doc_file:
        if current_user.is_authenticated:
            file_id = fs.put(doc_file, filename=doc_filename, user_id=str(current_user.get_id()))
    
    os.remove(temp_path)
    os.rmdir(temp_dir)

    quiz_url = url_for('get_doc', file_id=file_id, _external=True)
    return jsonify(success=True, quiz_url=quiz_url)

def get_file_stream(file_id):
    '''Retrieve a file from GridFS and return it as a stream.'''
    try:
        file_id = ObjectId(file_id)
        fs = GridFS(db)
        file = fs.get(file_id)
        file_stream = io.BytesIO(file.read())
        return file_stream
    except Exception as e:
        print(f"Error retrieving file: {e}")
        return None

def extract_text_from_pdf(file_stream):
    '''Extract text from a PDF file stream.'''
    file_stream.seek(0)
    reader = PdfReader(file_stream)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_pptx(file_stream):
    '''Extract text from a PPTX file stream.'''
    file_stream.seek(0)
    presentation = Presentation(file_stream)
    text = ""
    first_slide = presentation.slides[0]
    for shape in first_slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + "\n"
            break
    return text

def generate_questions(text):
    '''Generate quiz questions based on the given text.'''
    mcq_prompt = f"Create 20 multiple-choice questions based on the following text: \n{text}\nEach question should have 4 options (A, B, C, D), and indicate the correct answer."
    short_ans_prompt = f"Create 5 short answer questions based on the following text: \n{text}"
    long_ans_prompt = f"Create 2 long answer questions based on the following text: \n{text}"
    application_prompt = f"Create 1 application question based on the following text: \n{text}"

    mcq_text = call_openai_api(mcq_prompt)
    short_ans_text = call_openai_api(short_ans_prompt)
    long_ans_text = call_openai_api(long_ans_prompt)
    application_text = call_openai_api(application_prompt)

    questions = []
    if mcq_text:
        questions.extend(mcq_text.strip().split('\n\n'))
    if short_ans_text:
        questions.extend(short_ans_text.strip().split('\n\n'))
    if long_ans_text:
        questions.extend(long_ans_text.strip().split('\n\n'))
    if application_text:
        questions.extend(application_text.strip().split('\n\n'))

    return questions if questions else ["Failed to generate questions."]

def calculate_ratings():
    '''Calculate the star ratings and average rating based on all reviews.'''
    all_reviews = list(reviews_collection.find())
    star_counts = {1: 0, 2: 0, 3: 0, 4: 0, 5: 0}
    total_rating = 0
    total_count = 0

    for review in all_reviews:
        rating = review.get('star_rating')
        if rating is not None and rating != '' and rating != 0 and rating != []:
            rating = int(rating)
            star_counts[rating] += 1
            total_rating += rating
            total_count += 1

    average_rating = total_rating / total_count if total_count else 0
    return star_counts, average_rating

@app.route('/ratings')
@login_required
def ratings():
    '''Display the star ratings and average rating based on all reviews.'''
    star_counts, average_rating = calculate_ratings()
    return jsonify({'star_counts': star_counts, 'average_rating': average_rating})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
